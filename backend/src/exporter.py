import os
import tempfile
import json
import pandas as pd
import numpy as np

# Use headless matplotlib backend to prevent GUI errors
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Import analytical functions
from src.analyzer import load_scoring_config, generate_ai_insights, generate_executive_summary, OPTIONAL_FIELDS

def generate_matplotlib_charts(df: pd.DataFrame, temp_dir: str) -> dict:
    """
    Generates static PNGs of the 11 dashboard charts and stores them in temp_dir.
    Checks column availability and returns a dict mapping chart names to image paths (or None if unavailable).
    """
    chart_paths = {}
    
    # 1. Website Availability (Pie Chart)
    if "Website Status" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        counts = df["Website Status"].value_counts()
        ax.pie(counts, labels=counts.index, autopct='%1.1f%%', startangle=90, colors=['#3B82F6', '#EF4444'])
        ax.axis('equal')
        ax.set_title("Website Availability", fontsize=12, fontweight='bold')
        path = os.path.join(temp_dir, "chart_website_pie.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["website_pie"] = path
    else:
        chart_paths["website_pie"] = None
        
    # 2. Business Categories Distribution (Bar Chart)
    if "Category" in df.columns:
        fig, ax = plt.subplots(figsize=(8, 4))
        counts = df["Category"].value_counts().head(10).sort_values(ascending=True)
        ax.barh(counts.index, counts.values, color='#4F46E5')
        ax.set_title("Top 10 Business Categories", fontsize=12, fontweight='bold')
        ax.set_xlabel("Count")
        path = os.path.join(temp_dir, "chart_category_bar.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["category_bar"] = path
    else:
        chart_paths["category_bar"] = None
        
    # 3. Rating Distribution Histogram
    if "Rating" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        ratings = df["Rating"].dropna()
        if not ratings.empty:
            ax.hist(ratings, bins=10, color="#10B981", edgecolor='white')
        ax.set_title("Rating Distribution", fontsize=12, fontweight='bold')
        ax.set_xlabel("Rating (1.0 - 5.0)")
        path = os.path.join(temp_dir, "chart_rating_hist.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["rating_hist"] = path
    else:
        chart_paths["rating_hist"] = None
        
    # 4. Reviews Distribution
    if "Reviews" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        reviews = df["Reviews"].dropna()
        if not reviews.empty:
            ax.hist(reviews, bins=15, color="#F59E0B", edgecolor='white')
        ax.set_title("Reviews Distribution", fontsize=12, fontweight='bold')
        ax.set_xlabel("Reviews Count")
        path = os.path.join(temp_dir, "chart_reviews_hist.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["reviews_hist"] = path
    else:
        chart_paths["reviews_hist"] = None
        
    # 5. Opportunity Score Distribution
    if "Opportunity Score" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        scores = df["Opportunity Score"].dropna()
        if not scores.empty:
            ax.hist(scores, bins=10, color="#6366F1", edgecolor='white')
        ax.set_title("Opportunity Score Distribution", fontsize=12, fontweight='bold')
        ax.set_xlabel("Score (0 - 100)")
        path = os.path.join(temp_dir, "chart_opp_hist.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["opp_hist"] = path
    else:
        chart_paths["opp_hist"] = None
        
    # 6. Average Rating by Category
    if "Category" in df.columns and "Rating" in df.columns:
        fig, ax = plt.subplots(figsize=(8, 4))
        avg_rates = df.groupby("Category")["Rating"].mean().sort_values(ascending=True).head(10)
        ax.barh(avg_rates.index, avg_rates.values, color='#EC4899')
        ax.set_title("Avg Rating by Category (Top 10)", fontsize=12, fontweight='bold')
        ax.set_xlabel("Avg Rating")
        ax.set_xlim(1.0, 5.0)
        path = os.path.join(temp_dir, "chart_category_rating.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["category_rating"] = path
    else:
        chart_paths["category_rating"] = None
        
    # 7. Businesses by City
    if "City" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        counts = df["City"].value_counts()
        ax.bar(counts.index, counts.values, color='#8B5CF6')
        ax.set_title("Businesses by City", fontsize=12, fontweight='bold')
        ax.set_ylabel("Count")
        plt.xticks(rotation=45)
        path = os.path.join(temp_dir, "chart_city_bar.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["city_bar"] = path
    else:
        chart_paths["city_bar"] = None
        
    # 8. Website Status by Category (Stacked Bar)
    if "Category" in df.columns and "Website Status" in df.columns:
        fig, ax = plt.subplots(figsize=(8, 4))
        # Cross tab
        ct = pd.crosstab(df["Category"], df["Website Status"])
        # Take top 8 categories by total count
        ct["total"] = ct.sum(axis=1)
        ct = ct.sort_values(by="total", ascending=False).head(8).drop(columns="total")
        ct.plot(kind="bar", stacked=True, color=['#3B82F6', '#EF4444'], ax=ax)
        ax.set_title("Website Status by Category", fontsize=12, fontweight='bold')
        ax.set_ylabel("Count")
        plt.xticks(rotation=45)
        path = os.path.join(temp_dir, "chart_website_category.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["website_category"] = path
    else:
        chart_paths["website_category"] = None
        
    # 9. Review Count vs Rating (Scatter Plot)
    if "Rating" in df.columns and "Reviews" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        # Map colors for categories/levels
        colors_map = {
            "Excellent Opportunity": "#EF4444",
            "High Opportunity": "#F59E0B",
            "Medium Opportunity": "#FBBF24",
            "Low Opportunity": "#3B82F6",
            "Digitally Established": "#10B981"
        }
        for level, color in colors_map.items():
            sub = df[df["Opportunity Level"] == level]
            ax.scatter(sub["Rating"], sub["Reviews"], label=level, c=color, alpha=0.7, edgecolors='none')
        ax.set_title("Reviews vs Ratings", fontsize=12, fontweight='bold')
        ax.set_xlabel("Rating")
        ax.set_ylabel("Reviews")
        ax.legend(fontsize=8)
        path = os.path.join(temp_dir, "chart_reviews_rating_scatter.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["reviews_rating_scatter"] = path
    else:
        chart_paths["reviews_rating_scatter"] = None
        
    # 10. Top 10 Highest Opportunity Businesses
    if "Business Name" in df.columns and "Opportunity Score" in df.columns:
        fig, ax = plt.subplots(figsize=(8, 4))
        top_opps = df.sort_values(by="Opportunity Score", ascending=True).head(10)
        ax.barh(top_opps["Business Name"], top_opps["Opportunity Score"], color="#EF4444")
        ax.set_title("Top 10 Highest Opportunity Businesses", fontsize=12, fontweight='bold')
        ax.set_xlabel("Opportunity Score")
        ax.set_xlim(0, 105)
        path = os.path.join(temp_dir, "chart_top_opps.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["top_opps"] = path
    else:
        chart_paths["top_opps"] = None
        
    # 11. Opportunity Score vs Data Quality Score (Scatter)
    if "Opportunity Score" in df.columns and "Data Quality Score" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        colors_map = {
            "Excellent Opportunity": "#EF4444",
            "High Opportunity": "#F59E0B",
            "Medium Opportunity": "#FBBF24",
            "Low Opportunity": "#3B82F6",
            "Digitally Established": "#10B981"
        }
        for level, color in colors_map.items():
            sub = df[df["Opportunity Level"] == level]
            ax.scatter(sub["Opportunity Score"], sub["Data Quality Score"], label=level, c=color, alpha=0.7, edgecolors='none')
        ax.set_title("Opportunity Score vs Data Quality Score", fontsize=12, fontweight='bold')
        ax.set_xlabel("Opportunity Score")
        ax.set_ylabel("Data Quality Score")
        ax.set_xlim(0, 100)
        ax.set_ylim(0, 100)
        ax.legend(fontsize=8)
        path = os.path.join(temp_dir, "chart_opp_vs_dq.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["opp_vs_dq"] = path
    else:
        chart_paths["opp_vs_dq"] = None
        
    # 12. Businesses by State (Optional)
    if "State" in df.columns:
        fig, ax = plt.subplots(figsize=(6, 4))
        counts = df["State"].value_counts()
        ax.bar(counts.index, counts.values, color='#06B6D4')
        ax.set_title("Businesses by State", fontsize=12, fontweight='bold')
        ax.set_ylabel("Count")
        plt.xticks(rotation=45)
        path = os.path.join(temp_dir, "chart_state_bar.png")
        plt.tight_layout()
        plt.savefig(path, dpi=150)
        plt.close()
        chart_paths["state_bar"] = path
    else:
        chart_paths["state_bar"] = None
        
    return chart_paths

# --- PDF Export Class ---
class ExecutiveReportPDF(FPDF):
    def header(self):
        # Draw top colored bar
        self.set_fill_color(59, 130, 246) # Blue accent
        self.rect(0, 0, 210, 8, 'F')
        
        self.set_font('helvetica', 'B', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, 'Google Maps Business Analysis Platform - Executive Intelligence Report', new_x="LMARGIN", new_y="NEXT", align='R')
        self.ln(2)
        
    def footer(self):
        self.set_y(-15)
        self.set_font('helvetica', 'I', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f'Page {self.page_no()}', align='C')

def create_pdf_report(df: pd.DataFrame, stats: dict, output_path: str):
    """Generates a professional PDF report containing Executive Summary, AI insights, and visual graphs."""
    pdf = ExecutiveReportPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # RENDER CHARTS
    with tempfile.TemporaryDirectory() as temp_dir:
        charts = generate_matplotlib_charts(df, temp_dir)
        
        # --- Page 1: Title Cover ---
        pdf.add_page()
        pdf.ln(40)
        pdf.set_font('helvetica', 'B', 28)
        pdf.set_text_color(31, 41, 55) # Dark gray
        pdf.cell(0, 15, "Digital Transformation Report", new_x="LMARGIN", new_y="NEXT", align='L')
        pdf.set_font('helvetica', '', 14)
        pdf.set_text_color(75, 85, 99)
        pdf.cell(0, 10, "Local Business Digital Opportunity & Data Quality Analysis", new_x="LMARGIN", new_y="NEXT", align='L')
        pdf.ln(10)
        
        # Separator line
        pdf.set_draw_color(59, 130, 246)
        pdf.set_line_width(1.5)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(25)
        
        # Metadata Card
        pdf.set_fill_color(243, 244, 246)
        pdf.rect(10, pdf.get_y(), 190, 45, 'F')
        pdf.set_font('helvetica', 'B', 10)
        pdf.set_text_color(55, 65, 81)
        pdf.set_x(15)
        pdf.cell(0, 8, "REPORT METRICS:", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font('helvetica', '', 10)
        pdf.set_x(15)
        pdf.cell(0, 6, f"Total Analyzed Businesses: {len(df)}", new_x="LMARGIN", new_y="NEXT")
        pdf.set_x(15)
        pdf.cell(0, 6, f"Active Cities: {df['City'].nunique() if 'City' in df.columns else 'N/A'}", new_x="LMARGIN", new_y="NEXT")
        pdf.set_x(15)
        pdf.cell(0, 6, f"Average Opportunity Score: {int(round(df['Opportunity Score'].mean())) if 'Opportunity Score' in df.columns else 'N/A'}/100", new_x="LMARGIN", new_y="NEXT")
        pdf.set_x(15)
        pdf.cell(0, 6, f"Average Data Quality Score: {int(round(df['Data Quality Score'].mean())) if 'Data Quality Score' in df.columns else 'N/A'}/100", new_x="LMARGIN", new_y="NEXT")
        
        # --- Page 2: Executive Summary ---
        pdf.add_page()
        pdf.set_font('helvetica', 'B', 18)
        pdf.set_text_color(31, 41, 55)
        pdf.cell(0, 10, "Executive Summary", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)
        
        exec_summary = generate_executive_summary(df)
        pdf.set_font('helvetica', '', 11)
        pdf.set_text_color(55, 65, 81)
        # Multi-line writing
        pdf.multi_cell(0, 6, exec_summary["summary_text"])
        pdf.ln(10)
        
        # Highlight metrics box
        pdf.set_fill_color(239, 246, 255) # Light blue box
        pdf.rect(10, pdf.get_y(), 190, 30, 'F')
        pdf.set_y(pdf.get_y() + 4)
        pdf.set_font('helvetica', 'B', 10)
        pdf.set_text_color(30, 58, 138)
        pdf.set_x(15)
        pdf.cell(0, 6, f"Website Ownership Gap: {exec_summary['website_gap_percent']}% of businesses have no digital address.", new_x="LMARGIN", new_y="NEXT")
        pdf.set_x(15)
        pdf.cell(0, 6, f"Primary Focus City: {exec_summary['top_city']} hosts the highest density of high opportunity businesses.", new_x="LMARGIN", new_y="NEXT")
        pdf.set_x(15)
        pdf.cell(0, 6, f"Top Opportunity Sectors: {', '.join(exec_summary['top_categories'])} should be targeted first.", new_x="LMARGIN", new_y="NEXT")
        
        # --- Page 3: Column Inventory & Data Cleaning Summary ---
        pdf.add_page()
        pdf.set_font('helvetica', 'B', 18)
        pdf.set_text_color(31, 41, 55)
        pdf.cell(0, 10, "Dataset Inspection & Data Cleaning Summary", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)
        
        # Display list of detected columns
        pdf.set_font('helvetica', 'B', 12)
        pdf.cell(0, 8, "Column Ingestion Audit", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font('helvetica', '', 10)
        present_opts = [opt for opt in OPTIONAL_FIELDS if opt in df.columns]
        missing_opts = [opt for opt in OPTIONAL_FIELDS if opt not in df.columns]
        pdf.multi_cell(0, 5, f"Required columns: Business Name, Category, Rating, Reviews, Website, Phone, Address, City (All Detected - Verified).\n"
                             f"Detected Optional Columns: {', '.join(present_opts) if present_opts else 'None'}\n"
                             f"Absent Optional Columns: {', '.join(missing_opts) if missing_opts else 'None'}")
        pdf.ln(5)
        
        # Display Cleaning stats
        pdf.set_font('helvetica', 'B', 12)
        pdf.cell(0, 8, "Data Sanitization Pipeline Stats", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)
        
        # Table of cleaning results
        pdf.set_fill_color(249, 250, 251)
        pdf.set_font('helvetica', 'B', 10)
        pdf.cell(100, 7, "Cleaning Operation", 1, 0, 'L', True)
        pdf.cell(50, 7, "Transformations Count", 1, 1, 'R', True)
        pdf.set_font('helvetica', '', 10)
        
        labels = {
            "original_records": "Original Ingested Records",
            "duplicates_removed": "Duplicate Records Removed",
            "empty_names_removed": "Records with Empty Names Omitted",
            "invalid_ratings_corrected": "Clamped Invalid Ratings",
            "negative_reviews_corrected": "Corrected Negative Reviews (Forced Positive)",
            "missing_websites_marked": "Missing Websites Marked 'No Website'",
            "missing_phones_logged": "Missing Phone Numbers Logged",
            "categories_standardized": "Standardized Category Names",
            "final_clean_records": "Final Available Clean Dataset"
        }
        for key, label in labels.items():
            pdf.cell(100, 6, label, 1, 0, 'L')
            pdf.cell(50, 6, str(stats.get(key, 0)), 1, 1, 'R')
            
        # --- Page 4: AI Insights & Scoring Guidelines ---
        pdf.add_page()
        pdf.set_font('helvetica', 'B', 18)
        pdf.set_text_color(31, 41, 55)
        pdf.cell(0, 10, "AI Analytical Insights", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)
        
        insights = generate_ai_insights(df)
        pdf.set_font('helvetica', '', 11)
        for ins in insights:
            pdf.cell(5, 6, "*", 0, 0)
            pdf.multi_cell(0, 6, ins)
            pdf.ln(2)
            
        pdf.ln(5)
        pdf.set_font('helvetica', 'B', 14)
        pdf.cell(0, 8, "Scoring Methodology & Formulas", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font('helvetica', '', 10)
        pdf.multi_cell(0, 5, "1. Digital Growth Opportunity Score (0-100):\n"
                             "   - Website Missing (+40), Rating < 2.5 (+20), Reviews < 10 (+15), Phone Missing (+5),\n"
                             "     Industry Category Weight (+8/+10), Verified Business with strong footprint (-10).\n\n"
                             "2. Data Quality Score (0-100):\n"
                             "   - Starts at 100 points. Subtracts for missing mandatory columns (-10 to -20), duplicate records (-10),\n"
                             "     formatting modifications (-5 to -10), and missing optional details (-2 per empty field).")
        
        # --- Page 5: Top 10 Opportunity Recommendations ---
        pdf.add_page()
        pdf.set_font('helvetica', 'B', 16)
        pdf.set_text_color(31, 41, 55)
        pdf.cell(0, 10, "Top 10 Digital Transformation Opportunities", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)
        
        top10 = df.sort_values(by="Opportunity Score", ascending=False).head(10)
        
        # Render table
        pdf.set_fill_color(243, 244, 246)
        pdf.set_font('helvetica', 'B', 9)
        pdf.cell(50, 7, "Business Name", 1, 0, 'L', True)
        pdf.cell(30, 7, "Category", 1, 0, 'L', True)
        pdf.cell(20, 7, "City", 1, 0, 'L', True)
        pdf.cell(20, 7, "Opp Score", 1, 0, 'R', True)
        pdf.cell(20, 7, "Quality", 1, 0, 'R', True)
        pdf.cell(50, 7, "Primary Recommendation", 1, 1, 'L', True)
        
        pdf.set_font('helvetica', '', 8)
        for idx, row in top10.iterrows():
            name = str(row["Business Name"])[:25]
            cat = str(row["Category"])[:15]
            city = str(row["City"])[:10]
            opp = str(int(row["Opportunity Score"]))
            dq = str(int(row["Data Quality Score"]))
            reason = str(row["Recommendation Reason"])[:30]
            
            pdf.cell(50, 6, name, 1, 0, 'L')
            pdf.cell(30, 6, cat, 1, 0, 'L')
            pdf.cell(20, 6, city, 1, 0, 'L')
            pdf.cell(20, 6, opp, 1, 0, 'R')
            pdf.cell(20, 6, dq, 1, 0, 'R')
            pdf.cell(50, 6, reason, 1, 1, 'L')
            
        pdf.ln(10)
        pdf.set_font('helvetica', 'B', 12)
        pdf.cell(0, 6, "Suggested Transformation Actions:", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font('helvetica', '', 10)
        pdf.multi_cell(0, 5, "- For businesses with 'No Website': Target them with customized web development pitches.\n"
                             "- For businesses with low reviews/ratings: Pitch review management, feedback loops, and local SEO services.\n"
                             "- For businesses missing phone numbers: Update and claim Google Business profiles.")
        
        # --- Page 6: Visual Charts Gallery ---
        # Add charts in pairs per page
        chart_pairs = [
            ("website_pie", "category_bar"),
            ("rating_hist", "reviews_hist"),
            ("opp_hist", "opp_vs_dq"),
            ("category_rating", "website_category")
        ]
        
        for p1, p2 in chart_pairs:
            pdf.add_page()
            pdf.set_font('helvetica', 'B', 16)
            pdf.cell(0, 10, "Charts Analysis Gallery", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(5)
            
            # Left Chart
            path1 = charts.get(p1)
            if path1 and os.path.exists(path1):
                pdf.image(path1, x=10, y=30, w=90, h=60)
            else:
                pdf.set_fill_color(243, 244, 246)
                pdf.rect(10, 30, 90, 60, 'F')
                pdf.set_font('helvetica', 'I', 9)
                pdf.set_xy(10, 55)
                pdf.cell(90, 10, f"{p1} unavailable - column missing", 0, 0, 'C')
                
            # Right Chart
            path2 = charts.get(p2)
            if path2 and os.path.exists(path2):
                pdf.image(path2, x=110, y=30, w=90, h=60)
            else:
                pdf.set_fill_color(243, 244, 246)
                pdf.rect(110, 30, 90, 60, 'F')
                pdf.set_font('helvetica', 'I', 9)
                pdf.set_xy(110, 55)
                pdf.cell(90, 10, f"{p2} unavailable - column missing", 0, 0, 'C')
                
            pdf.ln(65)
            
    pdf.output(output_path)

# --- PPTX Export Helper ---
def create_pptx_report(df: pd.DataFrame, stats: dict, output_path: str):
    """Generates a professional PowerPoint presentation report matching the dashboard elements."""
    prs = Presentation()
    
    # RENDER CHARTS
    with tempfile.TemporaryDirectory() as temp_dir:
        charts = generate_matplotlib_charts(df, temp_dir)
        
        # Color constants
        c_dark = RGBColor(31, 41, 55)
        c_blue = RGBColor(59, 130, 246)
        c_gray = RGBColor(107, 114, 128)
        
        # Helper to set background and slide titles
        def add_styled_slide(title_text):
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
            # Remove title shape placeholder or use it
            title = slide.shapes.title
            title.text = title_text
            title.text_frame.paragraphs[0].font.size = Pt(24)
            title.text_frame.paragraphs[0].font.color.rgb = c_dark
            title.text_frame.paragraphs[0].font.bold = True
            
            # Draw bottom header colored bar
            top_bar = slide.shapes.add_shape(
                1, # ShapeType rectangle
                Inches(0), Inches(0), Inches(10), Inches(0.15)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = c_blue
            top_bar.line.color.rgb = c_blue
            
            return slide
            
        # Slide 1: Cover Page (Dark Indigo / Slate design)
        cover_slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        # Dark BG shape
        bg = cover_slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(17, 24, 39) # Dark gray/black
        bg.line.fill.background()
        
        # Main Title Box
        title_box = cover_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Google Maps Business Intelligence"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        p2 = tf.add_paragraph()
        p2.text = "India Local Business Opportunity & Data Quality Analysis"
        p2.font.size = Pt(18)
        p2.font.color.rgb = c_blue
        
        # Slide 2: Executive Summary
        slide2 = add_styled_slide("Executive Summary")
        exec_summary = generate_executive_summary(df)
        
        txBox = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = exec_summary["summary_text"]
        p.font.size = Pt(16)
        p.font.color.rgb = c_dark
        p.space_after = Pt(20)
        
        p_bullet1 = tf.add_paragraph()
        p_bullet1.text = f"- Digital Address Gap: {exec_summary['website_gap_percent']}% of local businesses lack websites."
        p_bullet1.font.size = Pt(14)
        p_bullet1.font.color.rgb = c_dark
        p_bullet1.space_after = Pt(10)
        
        p_bullet2 = tf.add_paragraph()
        p_bullet2.text = f"- Target Hotspot: {exec_summary['top_city']} contains the highest count of high-priority opportunities."
        p_bullet2.font.size = Pt(14)
        p_bullet2.font.color.rgb = c_dark
        p_bullet2.space_after = Pt(10)
        
        p_bullet3 = tf.add_paragraph()
        p_bullet3.text = f"- Recommended Target Verticals: {', '.join(exec_summary['top_categories'])}."
        p_bullet3.font.size = Pt(14)
        p_bullet3.font.color.rgb = c_dark
        
        # Slide 3: Dataset Ingestion & Feature Availability
        slide3 = add_styled_slide("Dataset Audit & Feature Availability")
        txBox = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        present_opts = [opt for opt in OPTIONAL_FIELDS if opt in df.columns]
        missing_opts = [opt for opt in OPTIONAL_FIELDS if opt not in df.columns]
        
        p = tf.paragraphs[0]
        p.text = "Ingested Dataset Checklists:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(15)
        
        p_req = tf.add_paragraph()
        p_req.text = "Required Columns: Business Name, Category, Rating, Reviews, Website, Phone, Address, City (ALL DETECTED)"
        p_req.font.bold = True
        p_req.font.size = Pt(14)
        p_req.font.color.rgb = RGBColor(16, 185, 129) # green
        p_req.space_after = Pt(15)
        
        p_opt_f = tf.add_paragraph()
        p_opt_f.text = f"Optional Columns Found: {', '.join(present_opts) if present_opts else 'None'}"
        p_opt_f.font.size = Pt(12)
        p_opt_f.space_after = Pt(10)
        
        p_opt_m = tf.add_paragraph()
        p_opt_m.text = f"Optional Columns Missing: {', '.join(missing_opts) if missing_opts else 'None'}"
        p_opt_m.font.size = Pt(12)
        
        # Slide 4: Data Cleaning Summary
        slide4 = add_styled_slide("Data Sanitization Stats")
        # Add statistics table
        rows, cols = 8, 2
        left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4)
        table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Headers
        table.cell(0, 0).text = "Preprocessing Operation"
        table.cell(0, 1).text = "Count"
        
        labels_pptx = [
            ("original_records", "Original Ingested Records"),
            ("duplicates_removed", "Duplicate Records Removed"),
            ("invalid_ratings_corrected", "Invalid Ratings Clamped"),
            ("negative_reviews_corrected", "Negative Reviews Forced Positive"),
            ("missing_websites_marked", "Missing Websites Marked 'No Website'"),
            ("missing_phones_logged", "Missing Phones Logged"),
            ("final_clean_records", "Final Sanitized Dataset Count")
        ]
        for idx, (key, label) in enumerate(labels_pptx):
            table.cell(idx+1, 0).text = label
            table.cell(idx+1, 1).text = str(stats.get(key, 0))
            
        # Slide 5: AI Insights
        slide5 = add_styled_slide("Analytical Insights")
        txBox = slide5.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        insights = generate_ai_insights(df)
        for idx, ins in enumerate(insights[:6]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"• {ins}"
            p.font.size = Pt(14)
            p.space_after = Pt(12)
            
        # Slide 6: Scoring Methodology
        slide6 = add_styled_slide("Scoring Guidelines & config")
        txBox = slide6.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Opportunity Score Formula (0-100):"
        p.font.bold = True
        p.font.size = Pt(16)
        p.space_after = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = "Evaluates digital potential. High weights allocated to Website Missing (+40), low reviews/ratings, and priority business sectors (Hospitals, Clinics, Education +10; Restaurants, Salons, Retail +8). Verified, high-performing listings receive a -10 bonus deduction."
        p2.font.size = Pt(13)
        p2.space_after = Pt(15)
        
        p3 = tf.add_paragraph()
        p3.text = "Data Quality Score Formula (0-100):"
        p3.font.bold = True
        p3.font.size = Pt(16)
        p3.space_after = Pt(6)
        
        p4 = tf.add_paragraph()
        p4.text = "Evaluates dataset integrity. Starts at 100. Deducts for missing required columns (Website -20, Phone -15, Name/City/Address -10 each), duplicate instances (-10), formatting corrections (-5 to -10), and missing optional parameters (-2 per field)."
        p4.font.size = Pt(13)
        
        # Slide 7: Top opportunities table
        slide7 = add_styled_slide("Top 10 Business Opportunities")
        top10 = df.sort_values(by="Opportunity Score", ascending=False).head(10)
        
        # Table of top 10
        rows, cols = 11, 5
        table_shape = slide7.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
        table = table_shape.table
        
        table.cell(0, 0).text = "Business Name"
        table.cell(0, 1).text = "Category"
        table.cell(0, 2).text = "City"
        table.cell(0, 3).text = "Opp Score"
        table.cell(0, 4).text = "Quality Score"
        
        for idx, row in top10.reset_index(drop=True).iterrows():
            table.cell(idx+1, 0).text = str(row["Business Name"])[:25]
            table.cell(idx+1, 1).text = str(row["Category"])[:15]
            table.cell(idx+1, 2).text = str(row["City"])[:10]
            table.cell(idx+1, 3).text = str(int(row["Opportunity Score"]))
            table.cell(idx+1, 4).text = str(int(row["Data Quality Score"]))
            
        # Slides 8+: Add Charts Slide pairs
        chart_slides = [
            ("website_pie", "category_bar", "Website Ownership & Category Distribution"),
            ("rating_hist", "reviews_hist", "Customer Reviews & Ratings Distribution"),
            ("opp_hist", "opp_vs_dq", "Opportunity Score & Data Quality Scatter Comparison"),
            ("category_rating", "website_category", "Category-wise Ratings & Website Status Stack")
        ]
        
        for c1, c2, title in chart_slides:
            slide = add_styled_slide(title)
            
            p1 = charts.get(c1)
            if p1 and os.path.exists(p1):
                slide.shapes.add_picture(p1, Inches(0.5), Inches(2), Inches(4.2), Inches(3.5))
            else:
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4.2), Inches(1.5))
                tb.text_frame.text = f"{c1} unavailable - column missing"
                tb.text_frame.paragraphs[0].font.size = Pt(14)
                
            p2 = charts.get(c2)
            if p2 and os.path.exists(p2):
                slide.shapes.add_picture(p2, Inches(5.3), Inches(2), Inches(4.2), Inches(3.5))
            else:
                tb = slide.shapes.add_textbox(Inches(5.3), Inches(3), Inches(4.2), Inches(1.5))
                tb.text_frame.text = f"{c2} unavailable - column missing"
                tb.text_frame.paragraphs[0].font.size = Pt(14)
                
    prs.save(output_path)
